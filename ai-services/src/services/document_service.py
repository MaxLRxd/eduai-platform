"""Extracción de texto de documentos (PDF, DOCX, PPTX, TXT)."""

import io

import docx
import pptx
import pypdf
import structlog

logger = structlog.get_logger(__name__)


class DocumentService:
    def extract_text(self, data: bytes, filename: str = "") -> str:
        ext = (filename.rsplit(".", 1)[-1] if "." in filename else "").lower()
        try:
            if ext == "pdf" or data[:4] == b"%PDF":
                return self._from_pdf(data)
            if ext in ("docx", "doc"):
                return self._from_docx(data)
            if ext in ("pptx", "ppt"):
                return self._from_pptx(data)
            if ext == "txt":
                return data.decode("utf-8", errors="ignore").strip()
            if data[:2] == b"PK":
                text = self._from_docx(data)
                if not text:
                    text = self._from_pptx(data)
                return text
            return data.decode("utf-8", errors="ignore").strip()
        except Exception as exc:
            logger.warning("document_extract_failed", filename=filename, error=str(exc))
            return ""

    def _from_pdf(self, data: bytes) -> str:
        reader = pypdf.PdfReader(io.BytesIO(data))
        pages = [(page.extract_text() or "") for page in reader.pages]
        return "\n".join(pages).strip()

    def _from_docx(self, data: bytes) -> str:
        document = docx.Document(io.BytesIO(data))
        paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
        return "\n".join(paragraphs).strip()

    def _from_pptx(self, data: bytes) -> str:
        presentation = pptx.Presentation(io.BytesIO(data))
        parts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        parts.append(text)
        return "\n".join(parts).strip()
